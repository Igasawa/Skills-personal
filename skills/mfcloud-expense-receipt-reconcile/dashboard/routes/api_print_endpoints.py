from __future__ import annotations

from datetime import datetime
from pathlib import Path
import hashlib
import re
import shutil
import subprocess
import sys
import threading
import zipfile
from uuid import uuid4

from fastapi import APIRouter, File, HTTPException, Request, UploadFile, Query
from fastapi.responses import FileResponse, JSONResponse

from services import core, core_scheduler, core_shared
from services import pptx_polish_portable

from .api_helpers import (
    _actor_from_request,
    _extract_print_file_paths,
    _merge_pdfs,
    _open_file,
    _safe_print_source,
    _source_list_path,
    _source_manifest_path,
)

PPTX_JOBS_ROOT_NAME = "pptx_polish"
PPTX_JOB_FOLDER = "jobs"
PPTX_JOB_ID_RE = r"^job_[0-9]{8}_[0-9]{6}_[a-f0-9]{6}$"
PPTX_MAX_HISTORY = 120


def _normalize_job_id(job_id: str) -> str:
    safe = str(job_id or "").strip()
    if len(safe) > 80:
        raise HTTPException(status_code=400, detail="Invalid job id.")
    if not re.match(PPTX_JOB_ID_RE, safe):
        raise HTTPException(status_code=400, detail="Invalid job id.")
    return safe


def _pptx_jobs_root() -> Path:
    return core._artifact_root() / PPTX_JOBS_ROOT_NAME / PPTX_JOB_FOLDER


def _pptx_job_path(job_id: str) -> Path:
    return _pptx_jobs_root() / _normalize_job_id(job_id)


def _new_pptx_job_id() -> str:
    return f"job_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid4().hex[:6]}"


def _prune_pptx_job_history() -> None:
    root = _pptx_jobs_root()
    if not root.exists():
        return
    all_jobs = [path.name for path in root.iterdir() if path.is_dir() and PPTX_JOB_ID_RE.match(path.name)]
    for extra in sorted(all_jobs, reverse=True)[PPTX_MAX_HISTORY:]:
        shutil.rmtree(root / extra, ignore_errors=True)


def _read_pptx_job_meta(job_id: str) -> dict:
    path = _pptx_job_path(job_id) / "meta.json"
    if not path.exists():
        raise HTTPException(status_code=404, detail="Job not found.")
    data = core._read_json(path)
    if not isinstance(data, dict):
        raise HTTPException(status_code=500, detail="Job metadata is invalid.")
    return data


def _read_pptx_job_meta_safe(job_id: str) -> dict | None:
    try:
        return _read_pptx_job_meta(job_id)
    except HTTPException:
        return None


def _write_pptx_job_meta(job_id: str, payload: dict) -> None:
    path = _pptx_job_path(job_id) / "meta.json"
    core._write_json(path, payload)


def _serialize_job_list_entry(job_id: str, meta: dict) -> dict:
    return {
        "job_id": job_id,
        "status": str(meta.get("status") or "unknown"),
        "created_at": str(meta.get("created_at") or ""),
        "updated_at": str(meta.get("updated_at") or ""),
        "attempt": int(meta.get("attempt") or 0),
        "actor": meta.get("actor") if isinstance(meta.get("actor"), dict) else {},
        "input": meta.get("input") if isinstance(meta.get("input"), dict) else {},
        "output": meta.get("output") if isinstance(meta.get("output"), dict) else {},
        "summary": meta.get("summary") if isinstance(meta.get("summary"), dict) else {},
        "diff": meta.get("diff") if isinstance(meta.get("diff"), dict) else {},
        "engine": str(meta.get("engine") or ""),
        "message": str(meta.get("message") or ""),
        "error": meta.get("error"),
        "progress": meta.get("progress") if isinstance(meta.get("progress"), dict) else {},
    }


def _collect_pptx_summary(path: Path) -> dict:
    if not path.exists():
        return {"exists": False, "path": str(path), "size": 0, "slides": 0, "zip_entries": 0, "slide_hashes": {}}

    summary = {
        "exists": True,
        "path": str(path),
        "size": path.stat().st_size,
        "slides": 0,
        "zip_entries": 0,
        "slide_hashes": {},
        "slide_names": [],
    }
    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            summary["zip_entries"] = len(names)
            slides = [name for name in names if name.startswith("ppt/slides/") and name.endswith(".xml")]
            slides.sort()
            summary["slide_names"] = slides
            summary["slides"] = len(slides)
            hashes = {}
            for name in slides:
                try:
                    payload = zf.read(name)
                except Exception:
                    continue
                hashes[name] = hashlib.sha256(payload).hexdigest()
            summary["slide_hashes"] = hashes
    except Exception:
        return summary
    return summary


def _build_pptx_diff(before: dict, after: dict) -> dict:
    before_hashes = before.get("slide_hashes") if isinstance(before.get("slide_hashes"), dict) else {}
    after_hashes = after.get("slide_hashes") if isinstance(after.get("slide_hashes"), dict) else {}
    before_slides = set(before_hashes.keys())
    after_slides = set(after_hashes.keys())
    added = sorted(after_slides - before_slides)
    removed = sorted(before_slides - after_slides)
    changed = 0
    unchanged = 0
    for key in before_slides & after_slides:
        if before_hashes.get(key) == after_hashes.get(key):
            unchanged += 1
        else:
            changed += 1
    before_size = int(before.get("size") or 0)
    after_size = int(after.get("size") or 0)
    return {
        "before_size": before_size,
        "after_size": after_size,
        "size_delta": after_size - before_size,
        "slides_before": int(before.get("slides") or 0),
        "slides_after": int(after.get("slides") or 0),
        "slide_hash_changes": changed,
        "slides_unchanged": unchanged,
        "slides_added": len(added),
        "slides_removed": len(removed),
        "added_slides": added[:20],
        "removed_slides": removed[:20],
    }


def _polish_with_python_pptx(source: Path, output: Path) -> tuple[bool, dict]:
    portable_status = ""
    try:
        portable_success, portable_details = pptx_polish_portable.polish_with_portable_skill(source, output)
        if portable_success:
            return True, portable_details
        if isinstance(portable_details, dict):
            portable_status = str(portable_details.get("reason") or "")
    except Exception as exc:
        portable_status = f"{type(exc).__name__}: {exc}"

    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt
    except Exception as exc:
        details = {
            "reason": f"python-pptx unavailable: {exc}",
            "engine": "fallback_copy",
        }
        if portable_status:
            details["portable_skill_reason"] = portable_status
        return False, details

    try:
        prs = Presentation(str(source))
        shape_processed = 0
        paragraph_processed = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                tf = getattr(shape, "text_frame", None)
                if tf is None:
                    continue
                for paragraph in tf.paragraphs:
                    runs = list(getattr(paragraph, "runs", []))
                    if not runs:
                        continue
                    paragraph_processed += 1
                    is_title = paragraph.level == 0
                    base_size = Pt(34) if is_title else Pt(18)
                    paragraph.font.name = paragraph.font.name or "Yu Gothic"
                    if not paragraph.font.bold:
                        paragraph.font.bold = True if is_title else False
                    if paragraph.font.size is None:
                        paragraph.font.size = base_size
                    if paragraph.alignment is None:
                        paragraph.alignment = PP_ALIGN.LEFT
                    for run in runs:
                        font = run.font
                        if not font.name:
                            font.name = "Yu Gothic"
                        if font.size is None:
                            font.size = base_size
                        shape_processed += 1

        prs.save(str(output))
        return (
            True,
            {
        "engine": "portable-pptx-skill",
                "shape_processed": shape_processed,
                "paragraph_processed": paragraph_processed,
            },
        )
    except Exception as exc:
        details = {
            "reason": f"python-pptx processing failed: {exc}",
            "engine": "python-pptx",
        }
        if portable_status:
            details["portable_skill_reason"] = portable_status
        return False, details


def _mark_job_status(
    *,
    job_id: str,
    status: str,
    message: str,
    progress: int = 0,
    attempt: int | None = None,
    error: str | None = None,
    engine: str | None = None,
    summary: dict | None = None,
    diff: dict | None = None,
    input_payload: dict | None = None,
    output_payload: dict | None = None,
) -> dict:
    meta = _read_pptx_job_meta(job_id)
    meta["status"] = status
    meta["message"] = str(message or "")
    meta["updated_at"] = datetime.now().isoformat(timespec="seconds")
    if error is None:
        meta.pop("error", None)
    else:
        meta["error"] = str(error)
    if attempt is not None:
        meta["attempt"] = int(attempt)
    if progress:
        meta["progress"] = {"value": int(progress)}
    elif "progress" not in meta:
        meta["progress"] = {"value": 0}
    if engine:
        meta["engine"] = str(engine)
    if summary is not None:
        meta["summary"] = summary
    if diff is not None:
        meta["diff"] = diff
    if input_payload is not None:
        meta["input"] = input_payload
    if output_payload is not None:
        meta["output"] = output_payload
    _write_pptx_job_meta(job_id, meta)
    return meta


def _run_pptx_polish_job(job_id: str, actor: dict | None = None) -> None:
    root = _pptx_job_path(job_id)
    input_path = root / "input.pptx"
    output_path = root / "polished.pptx"

    attempts = 1
    meta = _read_pptx_job_meta(job_id)
    attempts = int(meta.get("attempt") or 0) + 1

    _mark_job_status(
        job_id=job_id,
        status="running",
        message="PPTX polish job started.",
        progress=10,
        attempt=attempts,
        engine=meta.get("engine") if isinstance(meta.get("engine"), str) else "unknown",
        input_payload=meta.get("input") if isinstance(meta.get("input"), dict) else {},
        output_payload={"filename": output_path.name, "size": 0},
    )

    try:
        before_summary = _collect_pptx_summary(input_path)
        success, details = _polish_with_python_pptx(input_path, output_path)

        if not success:
            shutil.copy2(input_path, output_path)
            if not output_path.exists():
                raise RuntimeError("Failed to generate output file.")

        after_summary = _collect_pptx_summary(output_path)
        diff = _build_pptx_diff(before_summary, after_summary)
        output_payload = {
            "filename": output_path.name,
            "size": int(output_path.stat().st_size) if output_path.exists() else 0,
        }
        _mark_job_status(
            job_id=job_id,
            status="done",
            message=(
                "PPTX polish completed: "
                + (details.get("engine") if isinstance(details.get("engine"), str) else "fallback")
            ),
            progress=100,
            attempt=attempts,
            engine=str(details.get("engine") or "fallback"),
            summary={
                "before": before_summary,
                "after": after_summary,
                "details": details,
            },
            diff=diff,
            input_payload=meta.get("input") if isinstance(meta.get("input"), dict) else {},
            output_payload=output_payload,
        )
    except Exception as exc:
        output_payload = {
            "filename": output_path.name,
            "size": int(output_path.stat().st_size) if output_path.exists() else 0,
        }
        _mark_job_status(
            job_id=job_id,
            status="failed",
            message="PPTX polish failed.",
            progress=0,
            attempt=attempts,
            error=str(exc),
            engine=meta.get("engine") if isinstance(meta.get("engine"), str) else "unknown",
            input_payload=meta.get("input") if isinstance(meta.get("input"), dict) else {},
            output_payload=output_payload,
        )

def _iter_pptx_job_ids() -> list[str]:
    root = _pptx_jobs_root()
    if not root.exists():
        return []
    return [path.name for path in root.iterdir() if path.is_dir() and PPTX_JOB_ID_RE.match(path.name)]


def _scan_pptx_jobs(*, limit: int = PPTX_MAX_HISTORY) -> list[dict]:
    max_limit = min(limit if limit > 0 else 1, PPTX_MAX_HISTORY)
    jobs: list[dict] = []
    for job_id in sorted(_iter_pptx_job_ids(), reverse=True):
        if len(jobs) >= max_limit:
            break
        meta = _read_pptx_job_meta_safe(job_id)
        if not isinstance(meta, dict):
            continue
        jobs.append(_serialize_job_list_entry(job_id, meta))
    return jobs


def _require_downloadable_job(job_id: str, job: dict) -> Path:
    output_path = _pptx_job_path(job_id) / "polished.pptx"
    if not output_path.exists():
        raise HTTPException(status_code=409, detail="Output file is not ready.")
    if str(job.get("status") or "").strip().lower() != "done":
        raise HTTPException(status_code=409, detail="Job is not finished.")
    return output_path


def register_api_print_endpoints(router: APIRouter) -> None:
    @router.post("/api/print/{ym}/{source}")
    def api_print(ym: str, source: str, request: Request) -> JSONResponse:
        ym = core._safe_ym(ym)
        if source not in {"amazon", "rakuten"}:
            raise HTTPException(status_code=400, detail="Invalid source.")

        year, month = core._split_ym(ym)
        actor = _actor_from_request(request)
        try:
            core._assert_source_action_allowed(year, month, source, "print")
        except HTTPException as exc:
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_prepare",
                status="rejected",
                actor=actor,
                source=source,
                details={"reason": str(exc.detail)},
            )
            raise
        output_root = core._artifact_root() / ym
        scripts_dir = core.SKILL_ROOT / "scripts"
        reports_dir = output_root / "reports"
        exclude_orders_json = output_root / "reports" / "exclude_orders.json"
        print_script = reports_dir / "print_all.ps1"
        source_manifest_path = _source_manifest_path(reports_dir, source)
        source_list_path = _source_list_path(reports_dir, source)
        print_count: int | None = None

        cmd = [
            sys.executable,
            str(scripts_dir / "collect_print.py"),
            "--year",
            str(year),
            "--month",
            str(month),
            "--output-dir",
            str(output_root),
            "--sources",
            source,
            "--skip-shortcut-download",
        ]
        if exclude_orders_json.exists():
            cmd += ["--exclude-orders-json", str(exclude_orders_json)]

        try:
            res = subprocess.run(cmd, cwd=str(scripts_dir), capture_output=True, text=True, check=False)
            if res.returncode != 0:
                raise HTTPException(
                    status_code=500,
                    detail=(
                        "collect_print.py failed:\n"
                        f"cmd: {cmd}\n"
                        f"exit: {res.returncode}\n"
                        f"stdout:\n{res.stdout}\n"
                        f"stderr:\n{res.stderr}\n"
                    ),
                )

            manifest = core._read_json(source_manifest_path)
            if isinstance(manifest, dict):
                try:
                    print_count = int(manifest.get("count"))
                except Exception:
                    print_count = None

            if not source_manifest_path.exists():
                raise HTTPException(status_code=404, detail=f"{source_manifest_path.name} not found.")
            print_command = f"POST /api/print-run/{ym}/{source}"
            wf = core._read_workflow(reports_dir)
            section = wf.get(source) if isinstance(wf.get(source), dict) else {}
            section["print_prepared_at"] = datetime.now().isoformat(timespec="seconds")
            # Prepare is intentionally "before print": completion must be recorded manually.
            section.pop("printed_at", None)
            wf[source] = section
            core._write_workflow(reports_dir, wf)
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_prepare",
                status="success",
                actor=actor,
                source=source,
                details={
                    "print_script": str(print_script),
                    "print_manifest": str(source_manifest_path),
                    "print_list": str(source_list_path),
                    "count": print_count,
                },
            )

            return JSONResponse(
                {
                    "status": "ok",
                    "source": source,
                    "count": print_count,
                    "print_script": str(print_script),
                    "print_manifest": str(source_manifest_path),
                    "print_list": str(source_list_path),
                    "print_command": print_command,
                    "excluded_pdfs_url": f"/runs/{ym}/excluded-pdfs",
                }
            )
        except HTTPException as exc:
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_prepare",
                status="failed",
                actor=actor,
                source=source,
                details={
                    "reason": str(exc.detail),
                    "print_script": str(print_script),
                    "print_manifest": str(source_manifest_path),
                    "print_list": str(source_list_path),
                    "count": print_count,
                },
            )
            raise


    @router.post("/api/print/{ym}/{source}/complete")
    def api_print_complete(ym: str, source: str, request: Request) -> JSONResponse:
        ym = core._safe_ym(ym)
        if source not in {"amazon", "rakuten"}:
            raise HTTPException(status_code=400, detail="Invalid source.")

        year, month = core._split_ym(ym)
        actor = _actor_from_request(request)
        output_root = core._artifact_root() / ym
        reports_dir = output_root / "reports"
        source_manifest_path = _source_manifest_path(reports_dir, source)
        print_count: int | None = None

        try:
            core._assert_source_action_allowed(year, month, source, "print")
        except HTTPException as exc:
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_complete",
                status="rejected",
                actor=actor,
                source=source,
                details={"reason": str(exc.detail)},
            )
            raise

        try:
            wf = core._read_workflow(reports_dir)
            section = wf.get(source) if isinstance(wf.get(source), dict) else {}
            if not section.get("print_prepared_at"):
                raise HTTPException(
                    status_code=409,
                    detail=f"Print preparation is required before marking {source} print completion.",
                )
            if not source_manifest_path.exists():
                raise HTTPException(
                    status_code=404,
                    detail="Print preparation not found. Run print preparation first.",
                )

            manifest = core._read_json(source_manifest_path)
            if isinstance(manifest, dict):
                try:
                    print_count = int(manifest.get("count"))
                except Exception:
                    print_count = None

            section["printed_at"] = datetime.now().isoformat(timespec="seconds")
            wf[source] = section
            core._write_workflow(reports_dir, wf)
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_complete",
                status="success",
                actor=actor,
                source=source,
                details={"count": print_count},
            )

            return JSONResponse({"status": "ok", "source": source, "count": print_count})
        except HTTPException as exc:
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_complete",
                status="rejected" if exc.status_code in {400, 404, 409} else "failed",
                actor=actor,
                source=source,
                details={"reason": str(exc.detail), "count": print_count},
            )
            raise


    def _execute_source_print_run(ym: str, source: str, actor: dict[str, str]) -> JSONResponse:
        ym = core._safe_ym(ym)
        source = _safe_print_source(source)
        year, month = core._split_ym(ym)
        reports_dir = core._artifact_root() / ym / "reports"
        workflow = core._read_workflow(reports_dir)
        section = workflow.get(source) if isinstance(workflow.get(source), dict) else {}
        if not section.get("print_prepared_at"):
            detail = f"Print preparation is required before {source} bulk print run."
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_run",
                status="rejected",
                actor=actor,
                source=source,
                details={"reason": detail},
            )
            raise HTTPException(status_code=409, detail=detail)

        manifest_path = _source_manifest_path(reports_dir, source)
        manifest = core._read_json(manifest_path)
        if not manifest_path.exists() or not isinstance(manifest, dict):
            detail = f"{manifest_path.name} not found. Run print preparation first."
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_run",
                status="rejected",
                actor=actor,
                source=source,
                details={"reason": detail},
            )
            raise HTTPException(status_code=404, detail=detail)

        raw_paths = _extract_print_file_paths(manifest)
        if not raw_paths:
            detail = f"No print targets were found in {manifest_path.name}."
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_run",
                status="rejected",
                actor=actor,
                source=source,
                details={"reason": detail},
            )
            raise HTTPException(status_code=409, detail=detail)

        existing_paths: list[Path] = []
        missing_files: list[str] = []
        for raw_path in raw_paths:
            candidate = Path(raw_path)
            if candidate.exists() and candidate.is_file():
                existing_paths.append(candidate)
            else:
                missing_files.append(str(candidate))
        if not existing_paths:
            detail = (
                f"All target PDFs are missing for {source}. "
                f"manifest={manifest_path.name} missing_count={len(missing_files)}"
            )
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_run",
                status="failed",
                actor=actor,
                source=source,
                details={"reason": detail, "missing_count": len(missing_files)},
            )
            raise HTTPException(status_code=500, detail=detail)

        merged_pdf_path = reports_dir / f"print_merged_{source}.pdf"
        try:
            merged_count, merged_pages = _merge_pdfs(existing_paths, merged_pdf_path)
        except HTTPException as exc:
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_run",
                status="failed",
                actor=actor,
                source=source,
                details={"reason": str(exc.detail), "missing_count": len(missing_files)},
            )
            raise
        except Exception as exc:
            detail = f"Merged PDF generation failed: {exc}"
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_run",
                status="failed",
                actor=actor,
                source=source,
                details={"reason": detail, "missing_count": len(missing_files)},
            )
            raise HTTPException(status_code=500, detail=detail) from exc

        open_result = _open_file(merged_pdf_path)
        if open_result.returncode != 0:
            detail = (
                "Open merged PDF failed:\n"
                f"path: {merged_pdf_path}\n"
                f"exit: {open_result.returncode}\n"
                f"stdout:\n{open_result.stdout}\n"
                f"stderr:\n{open_result.stderr}\n"
            )
            core._append_audit_event(
                year=year,
                month=month,
                event_type="source_action",
                action="print_run",
                status="failed",
                actor=actor,
                source=source,
                details={"reason": detail, "missing_count": len(missing_files)},
            )
            raise HTTPException(status_code=500, detail=detail)

        core._append_audit_event(
            year=year,
            month=month,
            event_type="source_action",
            action="print_run",
            status="success",
            actor=actor,
            source=source,
            details={
                "mode": "manual_open",
                "count": merged_count,
                "merged_pages": merged_pages,
                "missing_count": len(missing_files),
                "missing_files": missing_files,
                "merged_pdf_path": str(merged_pdf_path),
            },
        )
        return JSONResponse(
            {
                "status": "ok",
                "ym": ym,
                "source": source,
                "print_mode": "manual_open",
                "count": merged_count,
                "missing_count": len(missing_files),
                "merged_pdf_path": str(merged_pdf_path),
            }
        )


    @router.post("/api/print-run/{ym}/{source}")
    def api_print_run_by_source(ym: str, source: str, request: Request) -> JSONResponse:
        actor = _actor_from_request(request)
        return _execute_source_print_run(ym, source, actor)


    @router.post("/api/print-run/{ym}")
    def api_print_run_legacy(ym: str, request: Request, source: str | None = None) -> JSONResponse:
        actor = _actor_from_request(request)
        if source:
            return _execute_source_print_run(ym, source, actor)
        ym = core._safe_ym(ym)
        year, month = core._split_ym(ym)
        detail = "Deprecated endpoint. Use /api/print-run/{ym}/{source}."
        core._append_audit_event(
            year=year,
            month=month,
            event_type="source_action",
            action="print_run",
            status="rejected",
            actor=actor,
            details={"reason": detail},
        )
        raise HTTPException(status_code=400, detail=detail)


    @router.post("/api/pptx/polish")
    def api_pptx_polish_upload(request: Request, file: UploadFile = File(...)) -> JSONResponse:
        actor = _actor_from_request(request)
        filename = str(file.filename or "").strip()
        if not filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="Only .pptx files are accepted.")
        job_id = _new_pptx_job_id()
        root = _pptx_job_path(job_id)
        root.mkdir(parents=True, exist_ok=True)
        input_path = root / "input.pptx"
        file.file.seek(0)
        with input_path.open("wb") as destination:
            shutil.copyfileobj(file.file, destination)
        size = input_path.stat().st_size
        now = datetime.now().isoformat(timespec="seconds")
        meta = {
            "status": "uploaded",
            "created_at": now,
            "updated_at": now,
            "attempt": 0,
            "actor": actor,
            "input": {"filename": filename, "size": size, "path": str(input_path)},
            "output": {"filename": "polished.pptx", "size": 0},
            "summary": {},
            "diff": {},
            "message": "Upload complete. Start polish to generate output.",
            "progress": {"value": 0},
            "engine": "python-pptx",
        }
        _write_pptx_job_meta(job_id, meta)
        _prune_pptx_job_history()
        return JSONResponse(
            {
                "status": "ok",
                "job_id": job_id,
                "run_api": f"/api/pptx/polish/{job_id}/run",
                "status_api": f"/api/pptx/polish/{job_id}",
            },
            headers={"Cache-Control": "no-store"},
        )


    @router.post("/api/pptx/polish/{job_id}/run")
    def api_pptx_polish_run(job_id: str, request: Request) -> JSONResponse:
        actor = _actor_from_request(request)
        job_id = _normalize_job_id(job_id)
        meta = _read_pptx_job_meta(job_id)
        status = str(meta.get("status") or "").strip().lower()
        if status == "running":
            return JSONResponse(
                {
                    "status": "ok",
                    "job_id": job_id,
                    "message": "Job is already running.",
                    "status_api": f"/api/pptx/polish/{job_id}",
                },
                headers={"Cache-Control": "no-store"},
            )
        if status not in {"uploaded", "failed", "done"}:
            raise HTTPException(status_code=409, detail="Job is not ready to run.")
        input_path = _pptx_job_path(job_id) / "input.pptx"
        if not input_path.exists():
            raise HTTPException(status_code=404, detail="Uploaded file not found.")
        meta["actor"] = actor
        _write_pptx_job_meta(job_id, meta)
        threading.Thread(target=_run_pptx_polish_job, args=(job_id, actor), daemon=True).start()
        running = _read_pptx_job_meta(job_id)
        return JSONResponse(
            {
                "status": "ok",
                "job_id": job_id,
                "message": "Polish started.",
                "job": running,
                "status_api": f"/api/pptx/polish/{job_id}",
            },
            headers={"Cache-Control": "no-store"},
        )


    @router.get("/api/pptx/polish/{job_id}")
    def api_pptx_polish_status(job_id: str) -> JSONResponse:
        job_id = _normalize_job_id(job_id)
        job = _read_pptx_job_meta(job_id)
        status = str(job.get("status") or "").strip().lower()
        if status == "done":
            output_path = _pptx_job_path(job_id) / "polished.pptx"
            if output_path.exists():
                output_payload = job.get("output") if isinstance(job.get("output"), dict) else {}
                output_payload = dict(output_payload)
                output_payload["size"] = int(output_path.stat().st_size)
                job["output"] = output_payload
                job["download_url"] = f"/api/pptx/polish/{job_id}/download"
            else:
                job["download_url"] = None
        else:
            job["download_url"] = None
        return JSONResponse({"status": "ok", "job": job}, headers={"Cache-Control": "no-store"})


    @router.get("/api/pptx/polish/{job_id}/download")
    def api_pptx_polish_download(job_id: str) -> FileResponse:
        job_id = _normalize_job_id(job_id)
        job = _read_pptx_job_meta(job_id)
        if str(job.get("status") or "").strip().lower() != "done":
            raise HTTPException(status_code=409, detail="Job is not completed.")
        output_path = _require_downloadable_job(job_id, job)
        return FileResponse(
            output_path,
            filename=output_path.name,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


    @router.get("/api/pptx/polish/jobs")
    def api_pptx_polish_jobs(limit: int = Query(default=20, ge=1, le=PPTX_MAX_HISTORY)) -> JSONResponse:
        jobs = _scan_pptx_jobs(limit=limit)
        for job in jobs:
            status = str(job.get("status") or "").strip().lower()
            job_id = str(job.get("job_id") or "")
            if status == "done":
                job["download_url"] = f"/api/pptx/polish/{job_id}/download"
            else:
                job["download_url"] = None
        return JSONResponse({"status": "ok", "jobs": jobs}, headers={"Cache-Control": "no-store"})


